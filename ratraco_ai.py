import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import sounddevice as sd
import numpy as np
import threading
import tempfile
import os
from scipy.io.wavfile import write
from docx import Document
from moviepy import VideoFileClip
from pydub import AudioSegment
import matplotlib.pyplot as plt
import networkx as nx
from langdetect import detect
from transformers import pipeline
import librosa
import time
from pydub.utils import make_chunks
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import win32com.client
import cv2
import pytesseract
from PIL import Image
import io
import pyttsx3
from gtts import gTTS
import pygame
from io import BytesIO
import seaborn as sns
import random
from ultralytics import YOLO
import torch
import numpy as np
from PIL import Image, ImageDraw
import re

# Configure Tesseract executable path
pytesseract.pytesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Load YOLOv8n model
try:
    yolo_model = YOLO('yolov8n.pt')
except:
    print("Warning: YOLOv8n model not found. Please download it first.")

def analyze_slide_content(slide_image):
    """Analyze slide content using YOLOv8n and OCR"""
    try:
        if slide_image is None:
            return None
            
        # Convert slide image to numpy array
        img_array = np.array(slide_image)
        
        # Run YOLOv8n detection
        results = yolo_model(img_array)
        
        # Initialize analysis results
        analysis = {
            'text': '',
            'charts': [],
            'images': [],
            'tables': [],
            'description': ''
        }
        
        # Process YOLO results
        for result in results:
            boxes = result.boxes
            for box in boxes:
                try:
                    cls = int(box.cls[0])
                    conf = float(box.conf[0])
                    x1, y1, x2, y2 = map(int, box.xyxy[0])
                    
                    # Get class name
                    class_name = result.names[cls]
                    
                    # Extract region of interest
                    roi = img_array[y1:y2, x1:x2]
                    
                    if class_name in ['chart', 'graph', 'plot']:
                        # Analyze chart
                        chart_analysis = analyze_chart(roi)
                        analysis['charts'].append(chart_analysis)
                        analysis['description'] += f"\nBiểu đồ: {chart_analysis['description']}"
                        
                    elif class_name in ['image', 'photo']:
                        # Analyze image
                        image_analysis = analyze_image(roi)
                        analysis['images'].append(image_analysis)
                        analysis['description'] += f"\nHình ảnh: {image_analysis['description']}"
                        
                    elif class_name in ['table']:
                        # Analyze table
                        table_analysis = analyze_table(roi)
                        analysis['tables'].append(table_analysis)
                        analysis['description'] += f"\nBảng: {table_analysis['description']}"
                    
                    # Extract text from the region
                    text = pytesseract.image_to_string(roi, lang='vie+eng')
                    if text.strip():
                        analysis['text'] += f"\n{text.strip()}"
                except Exception as e:
                    print(f"Error processing detection box: {str(e)}")
                    continue
        
        return analysis
        
    except Exception as e:
        print(f"Error in slide analysis: {str(e)}")
        return None

def analyze_chart(chart_image):
    """Analyze chart content"""
    try:
        # Convert to grayscale
        gray = cv2.cvtColor(chart_image, cv2.COLOR_BGR2GRAY)
        
        # Extract text
        text = pytesseract.image_to_string(gray, lang='vie+eng')
        
        # Detect chart type
        chart_type = detect_chart_type(chart_image)
        
        # Extract data points
        data_points = extract_chart_data(chart_image, chart_type)
        
        description = f"Biểu đồ {chart_type} "
        if data_points:
            description += f"với các giá trị: {', '.join(map(str, data_points))}"
        
        return {
            'type': chart_type,
            'data': data_points,
            'text': text,
            'description': description
        }
        
    except Exception as e:
        print(f"Error in chart analysis: {str(e)}")
        return {'type': 'unknown', 'description': 'Không thể phân tích biểu đồ'}

def analyze_image(image):
    """Analyze image content"""
    try:
        # Extract text
        text = pytesseract.image_to_string(image, lang='vie+eng')
        
        # Get image description
        description = "Hình ảnh"
        if text.strip():
            description += f" có nội dung: {text.strip()}"
            
        return {
            'text': text,
            'description': description
        }
        
    except Exception as e:
        print(f"Error in image analysis: {str(e)}")
        return {'description': 'Không thể phân tích hình ảnh'}

def analyze_table(table_image):
    """Analyze table content"""
    try:
        # Extract text
        text = pytesseract.image_to_string(table_image, lang='vie+eng')
        
        # Parse table structure
        rows = text.split('\n')
        table_data = []
        for row in rows:
            if row.strip():
                cells = re.split(r'\s{2,}', row.strip())
                table_data.append(cells)
        
        description = "Bảng dữ liệu"
        if table_data:
            description += f" với {len(table_data)} hàng và {len(table_data[0])} cột"
            
        return {
            'data': table_data,
            'text': text,
            'description': description
        }
        
    except Exception as e:
        print(f"Error in table analysis: {str(e)}")
        return {'description': 'Không thể phân tích bảng'}

def detect_chart_type(chart_image):
    """Detect type of chart"""
    try:
        # Convert to grayscale
        gray = cv2.cvtColor(chart_image, cv2.COLOR_BGR2GRAY)
        
        # Edge detection
        edges = cv2.Canny(gray, 50, 150)
        
        # Count lines and curves
        lines = cv2.HoughLinesP(edges, 1, np.pi/180, 50)
        circles = cv2.HoughCircles(gray, cv2.HOUGH_GRADIENT, 1, 20)
        
        if circles is not None:
            return "tròn"
        elif lines is not None and len(lines) > 5:
            return "cột"
        else:
            return "đường"
            
    except Exception as e:
        print(f"Error in chart type detection: {str(e)}")
        return "không xác định"

def extract_chart_data(chart_image, chart_type):
    """Extract data points from chart"""
    try:
        # Convert to grayscale
        gray = cv2.cvtColor(chart_image, cv2.COLOR_BGR2GRAY)
        
        # Extract text
        text = pytesseract.image_to_string(gray, lang='vie+eng')
        
        # Extract numbers
        numbers = re.findall(r'\d+(?:[.,]\d+)?', text)
        
        # Convert to float
        data_points = []
        for num in numbers:
            try:
                data_points.append(float(num.replace(',', '.')))
            except:
                continue
                
        return data_points
        
    except Exception as e:
        print(f"Error in data extraction: {str(e)}")
        return []

# Global variables
SAMPLE_RATE = 16000
recording_data = []
stream = None
is_recording = False
device_list = []

# AI models
asr_pipe = None
summarizer = None

pygame.mixer.init()
is_presentation_paused = False
current_audio = None

# Add these global variables at the top with other globals
user_voice_path = None
is_recording_voice = False
voice_recording_data = []

class RatracoAITool:
    def __init__(self, root):
        self.root = root
        self.setup_window()
        self.create_styles()
        self.create_widgets()
        self.load_microphone_devices()
        
        # Start loading AI models in background
        threading.Thread(target=self.load_models, daemon=True).start()
    
    def setup_window(self):
        """Configure the main application window"""
        self.root.title("🎙️ AI RATRACO - Ghi âm, Chuyển đổi, Tóm tắt & Mindmap")
        self.root.geometry("880x800")
        self.root.configure(bg="#f5f7fa")
        self.root.resizable(False, False)
        
        # Set application icon if available
        try:
            self.root.iconbitmap("ratraco_icon.ico")
        except:
            pass
    
    def create_styles(self):
        """Configure the ttk styles for the application"""
        self.style = ttk.Style()
        self.style.theme_use("clam")
        
        # Main styles
        self.style.configure("TFrame", background="#f5f7fa")
        self.style.configure("TButton", font=("Segoe UI", 10), padding=8)
        self.style.configure("TLabel", font=("Segoe UI", 11), background="#f5f7fa")
        self.style.configure("Title.TLabel", font=("Segoe UI", 18, "bold"), background="#f5f7fa")
        self.style.configure("Subtitle.TLabel", font=("Segoe UI", 12), background="#f5f7fa")
        self.style.configure("Footer.TLabel", font=("Segoe UI", 9), background="#f5f7fa", foreground="#666666")
        
        # Button styles
        self.style.configure("Record.TButton", background="#4CAF50", foreground="white")
        self.style.configure("Stop.TButton", background="#F44336", foreground="white")
        self.style.configure("Action.TButton", background="#2196F3", foreground="white")
        
        # Progress bar style
        self.style.configure("TProgressbar", thickness=8, background="#2196F3")
    
    def create_widgets(self):
        """Create and layout all UI elements"""
        # Container frames
        self.header_frame = ttk.Frame(self.root, style="TFrame")
        self.header_frame.pack(fill=tk.X, padx=20, pady=(20, 0))
        
        self.main_frame = ttk.Frame(self.root, style="TFrame")
        self.main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        self.footer_frame = ttk.Frame(self.root, style="TFrame")
        self.footer_frame.pack(fill=tk.X, padx=20, pady=(0, 10))
        
        # Header elements
        ttk.Label(self.header_frame, text="AI RATRACO", style="Title.TLabel").pack(side=tk.LEFT)
        
        # Status and progress
        self.status_frame = ttk.Frame(self.main_frame, style="TFrame")
        self.status_frame.pack(fill=tk.X, pady=10)
        
        self.status_label = ttk.Label(self.status_frame, text="⏳ Đang tải mô hình AI...", style="TLabel")
        self.status_label.pack(side=tk.TOP, anchor=tk.W)
        
        self.progress = ttk.Progressbar(self.status_frame, orient="horizontal", length=760, 
                                        mode="determinate", style="TProgressbar")
        self.progress.pack(side=tk.TOP, fill=tk.X, pady=5)
        
        # Microphone selection
        self.mic_frame = ttk.Frame(self.main_frame, style="TFrame")
        self.mic_frame.pack(fill=tk.X, pady=10)
        
        ttk.Label(self.mic_frame, text="🎤 Chọn Micro:", style="TLabel").pack(side=tk.LEFT)
        
        self.mic_var = tk.StringVar()
        self.mic_dropdown = ttk.Combobox(self.mic_frame, textvariable=self.mic_var, 
                                        state="readonly", width=65)
        self.mic_dropdown.pack(side=tk.LEFT, padx=10)
        
        ttk.Button(self.mic_frame, text="🔄 Refresh", 
                  command=self.load_microphone_devices, width=10).pack(side=tk.LEFT)
        
        # Buttons section
        self.buttons_frame = ttk.Frame(self.main_frame, style="TFrame")
        self.buttons_frame.pack(fill=tk.X, pady=20)
        
        # Recording control buttons
        self.rec_frame = ttk.LabelFrame(self.buttons_frame, text="Ghi âm", style="TFrame")
        self.rec_frame.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")
        
        ttk.Button(self.rec_frame, text="▶️ Bắt đầu ghi", 
                  command=self.start_recording, width=15).grid(row=0, column=0, padx=5, pady=5)
        ttk.Button(self.rec_frame, text="⏹️ Dừng & lưu", 
                  command=self.stop_recording, width=15).grid(row=0, column=1, padx=5, pady=5)
        
        # File conversion buttons
        self.convert_frame = ttk.LabelFrame(self.buttons_frame, text="Chuyển đổi", style="TFrame")
        self.convert_frame.grid(row=0, column=1, padx=10, pady=10, sticky="nsew")
        
        ttk.Button(self.convert_frame, text="🎵 Chuyển MP3/MP4", 
                  command=self.convert_file_to_text, width=20).grid(row=0, column=0, padx=5, pady=5)
        
        # Summarization buttons
        self.summary_frame = ttk.LabelFrame(self.buttons_frame, text="Tóm tắt & Mindmap", style="TFrame")
        self.summary_frame.grid(row=1, column=0, columnspan=2, padx=10, pady=10, sticky="nsew")
        
        ttk.Button(self.summary_frame, text="📄 Tóm tắt Word", 
                  command=self.summarize_docx, width=15).grid(row=0, column=0, padx=5, pady=5)
        ttk.Button(self.summary_frame, text="🧠 Vẽ mindmap", 
                  command=self.draw_mindmap, width=15).grid(row=0, column=1, padx=5, pady=5)
        
        # Modify PowerPoint section
        self.ppt_frame = ttk.LabelFrame(self.buttons_frame, text="PowerPoint", style="TFrame")
        self.ppt_frame.grid(row=2, column=0, columnspan=2, padx=10, pady=10, sticky="nsew")
        
        ttk.Button(self.ppt_frame, text="📊 Word sang PPTX", 
                  command=self.convert_word_to_pptx, width=15).grid(row=0, column=0, padx=5, pady=5)
        ttk.Button(self.ppt_frame, text="📝 Phân tích & Soạn kịch bản", 
                  command=self.analyze_presentation, width=20).grid(row=0, column=1, padx=5, pady=5)
        ttk.Button(self.ppt_frame, text="🎤 Bắt đầu thuyết trình", 
                  command=self.start_presentation, width=20, state='disabled').grid(row=0, column=2, padx=5, pady=5)
        self.pause_button = ttk.Button(self.ppt_frame, text="⏸️ Tạm dừng", 
                                     command=self.toggle_presentation, width=15, state='disabled')
        self.pause_button.grid(row=0, column=3, padx=5, pady=5)
        
        # Add voice recording button
        self.voice_button = ttk.Button(self.ppt_frame, text="🎙️ Ghi giọng nói", 
                                     command=self.toggle_voice_recording, width=15)
        self.voice_button.grid(row=0, column=4, padx=5, pady=5)
        
        # Add script preview frame
        self.script_frame = ttk.LabelFrame(self.buttons_frame, text="Kịch bản thuyết trình", style="TFrame")
        self.script_frame.grid(row=3, column=0, columnspan=2, padx=10, pady=10, sticky="nsew")
        
        # Add scrolled text widget for script
        self.script_text = scrolledtext.ScrolledText(self.script_frame, wrap=tk.WORD, width=80, height=10)
        self.script_text.grid(row=0, column=0, padx=5, pady=5, sticky="nsew")
        
        # Add save script button
        ttk.Button(self.script_frame, text="💾 Lưu kịch bản", 
                  command=self.save_script, width=15).grid(row=1, column=0, padx=5, pady=5)
        
        # Configure grid weights
        self.buttons_frame.columnconfigure(0, weight=1)
        self.buttons_frame.columnconfigure(1, weight=1)
        
        # Footer elements
        support_text = "Hỗ trợ SĐT/ZALO: 0985363602 / Email: kiniemboquenjerry@gmail.com"
        ttk.Label(self.footer_frame, text=support_text, style="Footer.TLabel").pack(side=tk.LEFT)
        
        copyright_text = "© 2025 RATRACO AI Tool | Developed with ❤️"
        ttk.Label(self.footer_frame, text=copyright_text, style="Footer.TLabel").pack(side=tk.RIGHT)
    
    # ------- Core functionality methods -------
    
    def load_microphone_devices(self):
        """Load and populate available microphone devices"""
        global device_list
        try:
            device_list = sd.query_devices()
            input_devices = [f"{i}: {dev['name']}" for i, dev in enumerate(device_list) 
                            if dev["max_input_channels"] > 0]
            
            self.mic_dropdown["values"] = input_devices
            if input_devices:
                self.mic_dropdown.current(0)
                self.status_label.config(text="✅ Đã tìm thấy thiết bị âm thanh")
            else:
                self.status_label.config(text="⚠️ Không tìm thấy thiết bị micro nào")
        except Exception as e:
            messagebox.showerror("Lỗi thiết bị", f"Không thể tải danh sách micro: {str(e)}")
    
    def load_models(self):
        """Load AI models in background thread"""
        global asr_pipe, summarizer
        
        for progress_value in range(0, 101, 10):
            self.progress["value"] = progress_value
            self.root.update_idletasks()
            time.sleep(0.1)
            
        try:
            # Load speech recognition model
            self.status_label.config(text="⏳ Đang tải mô hình nhận dạng giọng nói...")
            self.root.update_idletasks()
            asr_pipe = pipeline("automatic-speech-recognition", model="openai/whisper-small")
            
            # Load summarization model
            self.progress["value"] = 50
            self.status_label.config(text="⏳ Đang tải mô hình tóm tắt văn bản...")
            self.root.update_idletasks()
            summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")
            
            self.progress["value"] = 100
            self.status_label.config(text="✅ Đã sẵn sàng sử dụng AI RATRACO!")
        except Exception as e:
            self.status_label.config(text="❌ Lỗi tải mô hình AI!")
            messagebox.showerror("Lỗi tải mô hình AI", str(e))
    
    # ------- Recording methods -------
    
    def audio_callback(self, indata, frames, time, status):
        """Callback function for audio stream"""
        global recording_data, is_recording
        if is_recording:
            volume = np.linalg.norm(indata) * 10
            volume = min(volume, 100)
            self.progress["value"] = volume
            recording_data.append(indata.copy())
    
    def start_recording(self):
        """Start audio recording"""
        global is_recording, stream, recording_data, SAMPLE_RATE
        
        try:
            if not self.mic_var.get():
                messagebox.showwarning("Chưa chọn micro", "Vui lòng chọn thiết bị micro trước khi ghi âm.")
                return
                
            device_index = int(self.mic_var.get().split(":")[0])
            is_recording = True
            recording_data = []
            
            stream = sd.InputStream(
                device=device_index, 
                channels=1, 
                samplerate=SAMPLE_RATE, 
                callback=self.audio_callback
            )
            stream.start()
            
            self.status_label.config(text="🎙️ Đang ghi âm...")
            self.progress["value"] = 0
        except Exception as e:
            messagebox.showerror("Lỗi ghi âm", str(e))
    
    def stop_recording(self):
        """Stop audio recording and process the recorded data"""
        global is_recording, stream
        
        is_recording = False
        try:
            if stream:
                stream.stop()
                stream.close()
        except:
            pass
            
        self.status_label.config(text="🛠️ Đang xử lý AI...")
        self.progress["value"] = 0
        
        threading.Thread(target=self.process_recording, daemon=True).start()
    
    def process_recording(self):
        """Process the recorded audio data"""
        global recording_data, SAMPLE_RATE
        
        try:
            if not recording_data:
                self.status_label.config(text="⚠️ Không có dữ liệu ghi âm")
                return
                
            audio_data = np.concatenate(recording_data, axis=0)
            int16_audio = np.int16(audio_data * 32767)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_file:
                write(temp_file.name, SAMPLE_RATE, int16_audio)
                self.transcribe_with_ai(temp_file.name, "ket_qua_ghi_am.docx", show_progress=True)
                os.remove(temp_file.name)
        except Exception as e:
            messagebox.showerror("Lỗi xử lý", str(e))
            self.status_label.config(text="❌ Lỗi xử lý âm thanh.")
    
    # ------- File conversion methods -------
    
    def convert_file_to_text(self):
        """Convert audio file to text"""
        path = filedialog.askopenfilename(
            filetypes=[("Media files", "*.mp3 *.mp4")],
            title="Chọn file âm thanh/video cần chuyển đổi"
        )
        
        if not path: 
            return
            
        self.status_label.config(text="⏳ Đang chuyển đổi file âm thanh...")
        self.progress["value"] = 0
        
        threading.Thread(target=self.process_file_audio, args=(path,), daemon=True).start()
    
    def process_file_audio(self, path):
        """Process audio file for transcription"""
        try:
            # Convert MP3/MP4 to temporary WAV
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmpfile:
                if path.endswith(".mp4"):
                    clip = VideoFileClip(path)
                    clip.audio.write_audiofile(tmpfile.name, verbose=False, logger=None)
                else:
                    sound = AudioSegment.from_file(path)
                    sound.export(tmpfile.name, format="wav")

            # Load WAV file with pydub for chunking
            audio = AudioSegment.from_wav(tmpfile.name)
            chunk_length_ms = 30 * 1000  # 30 seconds chunks
            chunks = make_chunks(audio, chunk_length_ms)
            os.remove(tmpfile.name)

            full_text = ""
            total_chunks = len(chunks)

            # Create temp folder for chunks
            tmp_folder = tempfile.mkdtemp()

            for i, chunk in enumerate(chunks):
                chunk_name = os.path.join(tmp_folder, f"chunk{i}.wav")
                chunk.export(chunk_name, format="wav")

                self.status_label.config(text=f"⏳ Chuyển đổi đoạn {i+1}/{total_chunks}...")
                self.progress["value"] = int(((i + 1) / total_chunks) * 100)
                self.root.update()

                # Transcribe chunk
                text = asr_pipe(chunk_name)["text"]
                full_text += text + " "

                os.remove(chunk_name)

            # Remove temp folder
            os.rmdir(tmp_folder)

            # Save results
            self.save_to_word(full_text.strip(), "ket_qua_chuyen_file.docx")
            self.status_label.config(text="✅ Đã lưu ket_qua_chuyen_file.docx")
            messagebox.showinfo("Hoàn tất", "Đã chuyển đổi và lưu vào ket_qua_chuyen_file.docx")
            self.progress["value"] = 100

        except Exception as e:
            messagebox.showerror("Lỗi", str(e))
            self.status_label.config(text="❌ Lỗi chuyển file âm thanh.")
    
    # ------- Summarization methods -------
    
    def summarize_docx(self):
        """Summarize a Word document"""
        path = filedialog.askopenfilename(
            filetypes=[("Word files", "*.docx")], 
            title="Chọn file Word cần tóm tắt"
        )
        
        if not path:
            return
            
        self.status_label.config(text="⏳ Đang tóm tắt bằng AI...")
        self.progress["value"] = 0
        
        threading.Thread(target=self.process_summary, args=(path,), daemon=True).start()
    
    def process_summary(self, path):
        """Process document for summarization with key points extraction"""
        try:
            # Load and extract text from document
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            
            # Detect language to adjust summarization approach
            lang = "vi"
            try:
                lang = detect(text[:30])  # Use first 1000 chars for language detection(đang để nhận dạng theo độ dài 30 ký tự đầu để nhận diện ngôn ngữ)
            except:
                pass  # Default to Vietnamese if detection fails
                
            # Break text into manageable chunks for summarization(sẻ nhỏ file word để tóm tắt 512 frame)
            chunk_size = 512
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
            total_chunks = len(chunks)
            
            # Process each chunk
            summaries = []
            for idx, chunk in enumerate(chunks):
                if not chunk.strip():
                    continue
                    
                # Update progress
                self.progress["value"] = int(((idx+1)/total_chunks)*100)
                self.status_label.config(text=f"⏳ Tóm tắt phần ({idx+1}/{total_chunks})")
                self.root.update_idletasks()
                
                # Generate summary(cấu hình độ dài tóm tắt)
                try:
                    summary = summarizer(
                        chunk, 
                        #max_length=512, 
                        max_length=110, 
                        min_length=30, 
                        do_sample=False
                    )[0]["summary_text"]
                    
                    summaries.append(summary)
                except Exception as e:
                    print(f"Error generating summary for chunk {idx+1}: {str(e)}")
                    # Sử dụng nội dung gốc nếu không tạo được summary
                    summaries.append(chunk)
            
            # Extract key points from summaries
            key_points = self.extract_key_points(summaries, lang)
            
            # Format the final summary with key points
            final_summary = "TÓM TẮT CHÍNH:\n\n"
            final_summary += "\n".join(summaries)
            
            final_summary += "\n\nCÁC Ý CHÍNH:\n\n"
            for i, point in enumerate(key_points, 1):
                final_summary += f"{i}. {point}\n"
            
            # Save result
            self.save_to_word(final_summary, "tomtat.docx")
            self.status_label.config(text="✅ Tóm tắt xong!")
            messagebox.showinfo("Tóm tắt", "Đã lưu vào tomtat.docx")
            self.progress["value"] = 100
            
        except Exception as e:
            messagebox.showerror("Lỗi tóm tắt", str(e))
            self.status_label.config(text="❌ Lỗi tóm tắt.")
    
    def extract_key_points(self, summaries, lang="vi"):
        """Extract key points from the summarized text"""
        # Combine all summaries
        combined_text = " ".join(summaries)
        
        # Simple sentence-based extraction for key points
        sentences = []
        for line in combined_text.split("."):
            line = line.strip()
            if len(line) > 20:  # Ignore very short fragments
                sentences.append(line)
        
        # Filter sentences to find key points
        # For Vietnamese, use length and position heuristics
        key_points = []
        
        # Take first sentence from each summary as it often contains key information
        for summary in summaries:
            first_sentence = summary.split(".")[0].strip()
            if first_sentence and len(first_sentence) > 30 and first_sentence not in key_points:
                key_points.append(first_sentence)
        
        # Add other significant sentences
        for sentence in sentences:
            # Look for indicator phrases suggesting important points
            indicators = ["chính là", "quan trọng", "cần", "phải", "nên", "then chốt", "mấu chốt", 
                        "then chot", "mau chot", "đầu tiên", "kết luận", "tóm lại", "vậy là"]
                        
            is_key = False
            for indicator in indicators:
                if indicator in sentence.lower():
                    is_key = True
                    break
                    
            if is_key and sentence not in key_points and len(sentence) > 30:
                key_points.append(sentence)
        
        # If we have too few points, add more based on length
        if len(key_points) < 3:
            for sentence in sorted(sentences, key=len, reverse=True):
                if sentence not in key_points and len(sentence) > 40:
                    key_points.append(sentence)
                    if len(key_points) >= 5:
                        break
        
        # Return up to 5 key points
        return key_points[:5]
    
    # ------- Mindmap methods -------
    
    def draw_mindmap(self):
        """Create a mindmap from summarized text"""
        path = "tomtat.docx"
        
        if not os.path.exists(path):
            messagebox.showerror("Lỗi", "Không tìm thấy tomtat.docx. Vui lòng tóm tắt trước.")
            return
            
        self.status_label.config(text="⏳ Đang vẽ sơ đồ tư duy...")
        self.progress["value"] = 0
        
        threading.Thread(target=self.process_mindmap, args=(path,), daemon=True).start()
    
    def process_mindmap(self, path):
        """Generate mindmap visualization from document"""
        try:
            # Extract text from the document
            doc = Document(path)
            
            # Look for "CÁC Ý CHÍNH" section
            found_key_points = False
            key_points = []
            
            for p in doc.paragraphs:
                text = p.text.strip()
                
                if "CÁC Ý CHÍNH" in text:
                    found_key_points = True
                    continue
                    
                if found_key_points and text and text[0].isdigit() and "." in text[:3]:
                    # This is likely a numbered key point
                    point = text[text.find(".")+1:].strip()
                    if point:
                        key_points.append(point)
            
            # If no key points section found, extract from full text
            if not key_points:
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                topics = [s.strip(" .") for s in text.split(".") if len(s.strip()) > 20]
                key_points = topics[:7]  # Limit to 7 points for readability
            
            # Create network graph
            G = nx.Graph()
            
            # Create central node
            central_topic = "Tóm tắt"
            G.add_node(central_topic, size=2000, color="#FF5722")
            
            # Add key points as nodes
            for idx, point in enumerate(key_points):
                # Truncate long text
                node_text = point[:40] + "..." if len(point) > 40 else point
                node_id = f"Ý {idx+1}"
                
                G.add_node(node_id, size=1500, color="#2196F3", text=node_text)
                G.add_edge(central_topic, node_id, weight=3)
            
            # Create figure
            plt.figure(figsize=(14, 10), facecolor="#f8f9fa")
            
            # Position nodes
            pos = nx.spring_layout(G, k=0.5, iterations=50)
            
            # Draw nodes
            node_sizes = [data.get("size", 1000) for node, data in G.nodes(data=True)]
            node_colors = [data.get("color", "#1976D2") for node, data in G.nodes(data=True)]
            
            nx.draw_networkx_nodes(G, pos, node_size=node_sizes, node_color=node_colors, alpha=0.9)
            
            # Draw edges
            edge_weights = [G[u][v].get("weight", 1) for u, v in G.edges()]
            nx.draw_networkx_edges(G, pos, width=edge_weights, alpha=0.7, edge_color="#78909C")
            
            # Draw central node label
            central_pos = {central_topic: pos[central_topic]}
            nx.draw_networkx_labels(G, central_pos, {central_topic: central_topic}, 
                                   font_size=16, font_weight="bold", font_color="white")
            
            # Draw key point node labels
            point_labels = {}
            for node, data in G.nodes(data=True):
                if node != central_topic:
                    point_labels[node] = f"{node}\n{data.get('text', '')}"
            
            point_pos = {k: v for k, v in pos.items() if k in point_labels}
            nx.draw_networkx_labels(G, point_pos, point_labels, font_size=10, 
                                   font_weight="bold", font_color="black",
                                   bbox=dict(facecolor="white", alpha=0.7, boxstyle="round,pad=0.5"))
            
            plt.axis("off")
            plt.tight_layout()
            plt.savefig("mindmap.png", dpi=300, bbox_inches="tight")
            plt.close()
            
            self.status_label.config(text="🧠 Đã lưu sơ đồ tư duy (mindmap.png)")
            messagebox.showinfo("Sơ đồ tư duy", "Đã lưu sơ đồ tư duy vào mindmap.png")
            
            # Open the image
            try:
                os.startfile("mindmap.png")
            except:
                pass
                
        except Exception as e:
            messagebox.showerror("Lỗi", str(e))
            self.status_label.config(text="❌ Lỗi vẽ sơ đồ tư duy.")
    
    # ------- Utility methods -------
    
    def save_to_word(self, text, filename="ket_qua.docx"):
        """Save text to Word document"""
        doc = Document()
        
        # Add heading
        doc.add_heading('KẾT QUẢ XỬ LÝ AI RATRACO', 0)
        
        # Add timestamp
        doc.add_paragraph(f"Thời gian: {time.strftime('%d/%m/%Y %H:%M:%S')}")
        doc.add_paragraph("-------------")
        
        # Add content with proper formatting for sections
        if "TÓM TẮT CHÍNH:" in text:
            parts = text.split("TÓM TẮT CHÍNH:")
            doc.add_paragraph(parts[0])
            
            summary_parts = parts[1].split("CÁC Ý CHÍNH:")
            
            doc.add_heading('TÓM TẮT CHÍNH', level=1)
            doc.add_paragraph(summary_parts[0].strip())
            
            if len(summary_parts) > 1:
                doc.add_heading('CÁC Ý CHÍNH', level=1)
                
                # Process key points
                key_points = summary_parts[1].strip().split('\n')
                for point in key_points:
                    if point.strip():
                        doc.add_paragraph(point.strip(), style='List Bullet')
        else:
            # Simply add paragraphs
            paragraphs = text.split('\n')
            for p in paragraphs:
                if p.strip():
                    doc.add_paragraph(p.strip())
        
        # Add footer
        doc.add_paragraph("-------------")
        doc.add_paragraph("Tạo bởi công cụ AI RATRACO")
        
        # Save document
        doc.save(filename)
    
    def transcribe_with_ai(self, audio_path, save_as, show_progress=False):
        """Transcribe audio to text using AI"""
        try:
            if show_progress:
                self.progress["value"] = 0
                self.status_label.config(text="⏳ Đang chuyển đổi giọng nói thành văn bản...")
                self.root.update()
                
            text = asr_pipe(audio_path)["text"]
            self.save_to_word(text, save_as)
            
            self.status_label.config(text=f"✅ Đã lưu: {save_as}")
            self.progress["value"] = 100
            
            messagebox.showinfo("Hoàn tất", f"Đã lưu vào {save_as}")
            
            # Try to open the document
            try:
                os.startfile(save_as)
            except:
                pass
                
        except Exception as e:
            messagebox.showerror("Lỗi AI", str(e))
            self.status_label.config(text="❌ Lỗi chuyển đổi giọng nói.")

    def convert_word_to_pptx(self):
        """Convert Word document to PowerPoint presentation"""
        path = filedialog.askopenfilename(
            filetypes=[("Word files", "*.docx")],
            title="Chọn file Word cần chuyển đổi"
        )
        
        if not path:
            return
        
        self.status_label.config(text="⏳ Đang chuyển đổi Word sang PowerPoint...")
        self.progress["value"] = 0
        
        threading.Thread(target=self.process_word_to_pptx, args=(path,), daemon=True).start()

    def process_word_to_pptx(self, path):
        """Process Word document and convert to PowerPoint"""
        try:
            # Load Word document
            doc = Document(path)
            
            # Create new presentation
            prs = Presentation()
            
            # Set slide dimensions to 16:9
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            # Process each paragraph as a new slide
            for i, para in enumerate(doc.paragraphs):
                if not para.text.strip():
                    continue
                
                # Create new slide
                slide_layout = prs.slide_layouts[1]  # Use title and content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title = slide.shapes.title
                title.text = f"Slide {i+1}"
                
                # Add content
                content = slide.placeholders[1]
                content.text = para.text
                
                # Update progress
                self.progress["value"] = int(((i + 1) / len(doc.paragraphs)) * 100)
                self.status_label.config(text=f"⏳ Đang tạo slide {i+1}...")
                self.root.update_idletasks()
            
            # Save presentation
            output_path = os.path.splitext(path)[0] + "_presentation.pptx"
            prs.save(output_path)
            
            self.status_label.config(text=f"✅ Đã lưu: {output_path}")
            self.progress["value"] = 100
            
            messagebox.showinfo("Hoàn tất", f"Đã chuyển đổi và lưu vào {output_path}")
            
            # Try to open the presentation
            try:
                os.startfile(output_path)
            except:
                pass
            
        except Exception as e:
            messagebox.showerror("Lỗi", str(e))
            self.status_label.config(text="❌ Lỗi chuyển đổi Word sang PowerPoint.")

    def analyze_presentation(self):
        """Analyze presentation and generate detailed script"""
        path = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx")],
            title="Chọn file PowerPoint cần phân tích"
        )
        
        if not path:
            return
        
        self.status_label.config(text="⏳ Đang phân tích nội dung...")
        self.progress["value"] = 0
        
        threading.Thread(target=self.process_presentation_analysis, args=(path,), daemon=True).start()

    def process_presentation_analysis(self, path):
        """Process presentation and generate detailed script"""
        # Initialize variables for potential cleanup
        powerpoint_app = None # win32com PowerPoint application object
        # presentation object from python-pptx (prs) does not need explicit close

        try:
            print(f"[DEBUG] Starting analysis for file: {path}")
            # Load presentation using python-pptx
            # python-pptx can load .pptx files
            prs = Presentation(path)
            print("[DEBUG] Presentation loaded with python-pptx.")

            # Check if presentation was loaded successfully (python-pptx raises exception on failure)
            # We can check if there are slides as a basic check
            if not prs.slides:
                messagebox.showerror("Lỗi", "Không thể mở file PowerPoint hoặc file không có slide nào.")
                # Set status to idle
                self.status_label.config(text="✅ Đã sẵn sàng sử dụng AI RATRACO!")
                self.progress["value"] = 0
                # Ensure presentation button is disabled if no slides are found
                for widget in self.ppt_frame.winfo_children():
                    if isinstance(widget, ttk.Button) and "Bắt đầu thuyết trình" in widget.cget("text"):
                        widget.config(state='disabled')
                print("[DEBUG] No slides found or failed to load.")
                return

            total_slides = len(prs.slides)
            print(f"[DEBUG] Found {total_slides} slides.")

            # Store the path for later use in presentation mode
            self.current_presentation = path

            # Generate script content string
            script_content = "KỊCH BẢN THUYẾT TRÌNH\n\n"
            print("[DEBUG] Initial script_content created.")

            # We don't open the slideshow here, only during the actual presentation

            for i, slide in enumerate(prs.slides): # i is 0-based index
                print(f"[DEBUG] Processing slide {i+1}/{total_slides}...")
                # Update progress
                self.progress["value"] = int(((i + 1) / total_slides) * 100)
                self.status_label.config(text=f"⏳ Đang phân tích slide {i+1}/{total_slides}...")
                self.root.update_idletasks()

                # --- Step 1: Extract text from slide using python-pptx ---
                slide_text = ""
                extracted_text_parts = []
                try:
                    for shape in slide.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                             # Check if shape has text_frame attribute and it's not None
                             if shape.text_frame:
                                 # Use .text property which gets all text from paragraphs
                                 text_content = shape.text_frame.text.strip()
                                 if text_content:
                                     extracted_text_parts.append(text_content)
                                     # print(f"[DEBUG] Extracted text from shape: {text_content[:50]}...") # Log snippet
                    slide_text = "\n".join(extracted_text_parts) # Join text from all shapes
                    print(f"[DEBUG] Combined extracted text for slide {i+1}:\n{slide_text[:200]}...") # Log combined text snippet
                except Exception as text_e:
                    print(f"[ERROR] Error extracting text from slide {i+1}: {text_e}")
                    slide_text = "(Lỗi trích xuất văn bản)"


                # --- Step 2: Convert slide to image for analysis (using win32com) ---
                # Note: This can be resource intensive and potentially cause issues.
                # We will handle exceptions and continue without image analysis if it fails.
                slide_image = None
                temp_image_path = None # Define outside try for cleanup
                try:
                    print(f"[DEBUG] Attempting to convert slide {i+1} to image...")
                    # Need to use win32com to convert slide to image
                    # Create a new PowerPoint application instance
                    powerpoint_app = win32com.client.Dispatch("PowerPoint.Application")
                    print("[DEBUG] PowerPoint application instance created.")

                    # Create a temporary presentation to hold the single slide for export
                    # Using a more unique filename to avoid conflicts
                    temp_prs_path = os.path.join(tempfile.gettempdir(), f"temp_analysis_export_{os.path.basename(path)}_{i}_{os.getpid()}.pptx")
                    temp_prs = powerpoint_app.Presentations.Add()
                    print("[DEBUG] Temporary presentation created.")

                    # Copy the slide from the original presentation (prs) to the temporary one
                    # win32com uses 1-based indexing for slides, python-pptx enumerate is 0-based
                    # We need to copy slide i from prs (python-pptx object)
                    # win32com Paste method adds to the beginning (index 1)
                    try:
                        slide.Copy()
                        temp_prs.Slides.Paste()
                        print("[DEBUG] Slide copied and pasted into temporary presentation.")
                    except Exception as copy_paste_e:
                         print(f"[ERROR] Error during slide copy/paste for slide {i+1}: {copy_paste_e}")
                         # Attempt to close PowerPoint and temp_prs before re-raising or handling
                         try:
                             temp_prs.Close() if temp_prs else None
                             powerpoint_app.Quit() if powerpoint_app else None
                         except:
                              pass
                         raise copy_paste_e # Re-raise to be caught by main try block

                    # Save the temp presentation - crucial for win32com to export correctly
                    temp_prs.Save(temp_prs_path)
                    print(f"[DEBUG] Temporary presentation saved to {temp_prs_path}")

                    # Now open the saved temp presentation to export the slide image
                    # Re-opening can sometimes help with COM object stability
                    # Use different variable name to ensure the first temp_prs is not used accidentally
                    temp_prs_for_export = powerpoint_app.Presentations.Open(temp_prs_path, ReadOnly=True)
                    print(f"[DEBUG] Temporary presentation re-opened from {temp_prs_path}")

                    # Export the first (and only) slide in the temporary presentation
                    temp_image_path = os.path.join(tempfile.gettempdir(), f"slide_analysis_{i}_{os.getpid()}.png") # Unique image name
                    # Export the slide at index 1 (the pasted slide) in the temp presentation
                    temp_prs_for_export.Slides[1].Export(temp_image_path, "PNG")
                    print(f"[DEBUG] Slide exported to image: {temp_image_path}")

                    # Close the temporary presentations and quit PowerPoint
                    temp_prs_for_export.Close()
                    temp_prs.Close()
                    powerpoint_app.Quit()
                    print("[DEBUG] Temporary presentations closed and PowerPoint quit.")

                    # Clean up temporary presentation file
                    if os.path.exists(temp_prs_path):
                        try:
                            os.remove(temp_prs_path)
                            print(f"[DEBUG] Cleaned up temp presentation file: {temp_prs_path}")
                        except Exception as remove_e:
                             print(f"[ERROR] Error removing temp presentation file {temp_prs_path}: {remove_e}")

                    # Open the created image file using PIL
                    if os.path.exists(temp_image_path):
                         slide_image = Image.open(temp_image_path)
                         print(f"[DEBUG] Opened temp image file: {temp_image_path}")
                    else:
                         print(f"[ERROR] Temp image file not found after export: {temp_image_path}")

                except Exception as img_e:
                    print(f"[ERROR] Error converting slide {i+1} to image: {img_e}")
                    # Clean up any potential lingering PowerPoint instances if an error occurred
                    try:
                        if 'powerpoint_app' in locals() and powerpoint_app is not None:
                             print("[DEBUG] Attempting to quit PowerPoint after image conversion error.")
                             powerpoint_app.Quit()
                    except:
                        pass
                    slide_image = None # Ensure slide_image is None if conversion fails
                    # If image conversion fails, also skip analysis cleanup of the image file
                    temp_image_path = None

                # --- Step 3: Process slide content (using YOLO and OCR on image if image is available) ---
                analysis = None
                if slide_image:
                    print(f"[DEBUG] Starting analysis for slide {i+1} image...")
                    try:
                        # Pass extracted text as well, might help analysis context
                        # analysis will be a dict if successful, None otherwise
                        analysis = self.process_slide_content(slide_text, slide_image)
                        print(f"[DEBUG] Analysis completed for slide {i+1}. Result keys: {analysis.keys() if analysis else 'None'}")
                    except Exception as analysis_e:
                        print(f"[ERROR] Error during analysis of slide {i+1}: {analysis_e}")
                        # Continue even if analysis fails
                    finally:
                        # Clean up the temporary image file if it was created and opened by PIL
                        if slide_image:
                             try:
                                 # Close the PIL image before attempting to remove the file
                                 slide_image.close()
                                 print(f"[DEBUG] Closed PIL image for slide {i+1}.")
                             except Exception as close_e:
                                 print(f"[ERROR] Error closing PIL image for slide {i+1}: {close_e}")

                        if temp_image_path and os.path.exists(temp_image_path):
                            try:
                                os.remove(temp_image_path)
                                print(f"[DEBUG] Cleaned up temp image file: {temp_image_path}")
                            except Exception as cleanup_e:
                                print(f"[ERROR] Error cleaning up temp image file {temp_image_path}: {cleanup_e}")


                # --- Step 4: Build the content for the script for this slide ---
                slide_script_content_parts = []

                # Always add extracted text if available
                if slide_text.strip():
                     slide_script_content_parts.append("Văn bản slide:\n" + slide_text.strip())
                     print(f"[DEBUG] Added extracted text to parts for slide {i+1}.")

                # Add descriptions from visual analysis if analysis was successful
                visual_descriptions = []
                if analysis:
                    print(f"[DEBUG] Processing analysis results for slide {i+1}.")
                    # Check charts
                    if analysis.get('charts'):
                        print(f"[DEBUG] Found {len(analysis['charts'])} charts in slide {i+1}.")
                        for chart in analysis['charts']:
                             desc = f"Biểu đồ {chart.get('type', 'không xác định')}: {chart.get('description', 'không có mô tả chi tiết')}"
                             visual_descriptions.append(desc)
                             # Optionally add chart data if available
                             if chart.get('data'):
                                 visual_descriptions.append(f"Dữ liệu biểu đồ: {', '.join(map(str, chart['data']))}")

                    # Check images
                    if analysis.get('images'):
                        print(f"[DEBUG] Found {len(analysis['images'])} images in slide {i+1}.")
                        for image in analysis['images']:
                             desc = f"Hình ảnh: {image.get('description', 'không có mô tả chi tiết')}"
                             visual_descriptions.append(desc)
                             if image.get('text', '').strip():
                                 visual_descriptions.append(f"Văn bản trong ảnh: {image['text'].strip()}")

                    # Check tables
                    if analysis.get('tables'):
                        print(f"[DEBUG] Found {len(analysis['tables'])} tables in slide {i+1}.")
                        for table in analysis['tables']:
                             desc = f"Bảng dữ liệu: {table.get('description', 'không có mô tả chi tiết')}"
                             visual_descriptions.append(desc)
                             if table.get('data'):
                                 table_data_str = "; ".join([" | ".join(map(str, row)) for row in table['data']])
                                 visual_descriptions.append(f"Dữ liệu bảng: {table_data_str}")

                    # Add visual descriptions to slide_script_content_parts
                    # Add all descriptions from analysis, regardless of redundancy with original text for the script display.
                    # Redundancy can be handled during the speaking phase if needed.
                    for desc in visual_descriptions:
                        if desc.strip(): # Only add non-empty descriptions
                            slide_script_content_parts.append(desc.strip())
                    print(f"[DEBUG] Added {len(visual_descriptions)} visual descriptions to parts for slide {i+1}.")
                else:
                     print(f"[DEBUG] No analysis results for slide {i+1}.")


                # Combine all parts for the slide's script entry
                # Use two newlines to separate different sections (text, charts, images, tables) for readability
                # Filter out any empty parts before joining.
                combined_slide_script = "\n\n".join([part for part in slide_script_content_parts if part.strip()]).strip()
                print(f"[DEBUG] Combined script for slide {i+1}:\n{combined_slide_script[:300]}...")

                # --- Step 5: Add to main script variable ---
                script_content += f"\n--- Slide {i+1} ---\n"
                if combined_slide_script:
                    # Add the combined text if it's not empty
                    script_content += combined_slide_script + "\n"
                else:
                    # If no combined text, add a placeholder. This covers cases where both text extraction and analysis failed or found nothing.
                    script_content += "(Không có nội dung văn bản hoặc nội dung phân tích từ hình ảnh/bảng biểu)\n"

                print(f"[DEBUG] Finished processing slide {i+1}. Total script_content length so far: {len(script_content)}")

                # We remove the automatic speech playback and slide advancement during analysis.
                # These actions belong in the process_ai_presentation function which is triggered by the user.

            # --- Step 6: Display script in text widget ---
            print("[DEBUG] Analysis complete. Displaying script.")
            self.script_text.delete(1.0, tk.END)
            self.script_text.insert(tk.END, script_content)
            print("[DEBUG] Script displayed in text widget.")

            # --- Step 7: Enable presentation button and update status ---
            print("[DEBUG] Enabling presentation button.")
            for widget in self.ppt_frame.winfo_children():
                # Assuming the "Bắt đầu thuyết trình" button is the one to enable
                if isinstance(widget, ttk.Button) and "Bắt đầu thuyết trình" in widget.cget("text"):
                    widget.config(state='normal')

            # Disable pause button (as presentation hasn't started yet)
            self.pause_button.config(state='disabled')
            print("[DEBUG] Pause button disabled.")

            self.status_label.config(text="✅ Đã hoàn thành phân tích!")
            self.progress["value"] = 100
            print("[DEBUG] Status updated to complete.")

            # Show message to the user
            messagebox.showinfo("Hoàn tất", "Đã phân tích xong. Bạn có thể xem và chỉnh sửa kịch bản thuyết trình trước khi bắt đầu.")
            print("[DEBUG] Showing completion message box.")

        except Exception as e:
            # Catch any other unexpected errors during the analysis process
            print(f"[ERROR] An unexpected error occurred during analysis: {e}")
            messagebox.showerror("Lỗi", f"Lỗi trong quá trình phân tích: {str(e)}")
            self.status_label.config(text="❌ Lỗi phân tích.")
            self.progress["value"] = 0 # Reset progress on error
            self.pause_button.config(state='disabled')
            # Also disable the presentation button on error
            for widget in self.ppt_frame.winfo_children():
                if isinstance(widget, ttk.Button) and "Bắt đầu thuyết trình" in widget.cget("text"):
                    widget.config(state='disabled')

            # Clean up potential lingering win32com objects in case of error
            try:
                # Check if powerpoint_app was created and is still running
                if 'powerpoint_app' in locals() and powerpoint_app is not None:
                    print("[DEBUG] Attempting to quit PowerPoint app in error cleanup.")
                    powerpoint_app.Quit()
            except Exception as cleanup_e:
                print(f"[ERROR] Error during PowerPoint cleanup in error handling: {cleanup_e}")

            # No need to close python-pptx presentation object explicitly

        print("[DEBUG] Exiting process_presentation_analysis.")

    def slide_to_image(self, slide):
        """Convert PowerPoint slide to PIL Image"""
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                # Save slide as image using win32com
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = True
                
                # Create temporary presentation
                temp_prs = powerpoint.Presentations.Add()
                temp_slide = temp_prs.Slides.Add(1, 12)  # 12 = ppLayoutBlank
                
                # Copy content from original slide
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        textbox = temp_slide.Shapes.AddTextbox(1, shape.left, shape.top, shape.width, shape.height)
                        textbox.TextFrame.TextRange.Text = shape.text
                    elif hasattr(shape, "image"):
                        # Handle images
                        temp_slide.Shapes.AddPicture(shape.image.filename, False, True, shape.left, shape.top, shape.width, shape.height)
                
                # Export as image
                temp_slide.Export(tmp.name, "PNG")
                
                # Clean up
                temp_prs.Close()
                powerpoint.Quit()
                
                # Open image
                image = Image.open(tmp.name)
                
                # Clean up temp file
                os.unlink(tmp.name)
                
                return image
                
        except Exception as e:
            print(f"Error converting slide to image: {str(e)}")
            return None

    def save_script(self):
        """Save the current script to a file"""
        script = self.script_text.get(1.0, tk.END)
        path = filedialog.asksaveasfilename(
            defaultextension=".txt",
            filetypes=[("Text files", "*.txt")],
            title="Lưu kịch bản thuyết trình"
        )
        
        if path:
            try:
                with open(path, 'w', encoding='utf-8') as f:
                    f.write(script)
                messagebox.showinfo("Thành công", "Đã lưu kịch bản!")
            except Exception as e:
                messagebox.showerror("Lỗi", f"Không thể lưu file: {str(e)}")

    def start_presentation(self):
        """Start the presentation with the current script"""
        if not hasattr(self, 'current_presentation'):
            messagebox.showerror("Lỗi", "Vui lòng phân tích presentation trước!")
            return
        
        # Get the current script
        script = self.script_text.get(1.0, tk.END)
        
        # Start presentation in a new thread
        threading.Thread(target=self.process_ai_presentation, 
                        args=(self.current_presentation, script), daemon=True).start()

    def process_ai_presentation(self, path, script):
        """Process PowerPoint file for AI presentation with custom script"""
        powerpoint = None
        presentation = None
        slideshow_window = None

        try:
            # Enable pause button
            self.pause_button.config(state='normal')

            # Initialize PowerPoint application
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True

            # Open presentation
            # Use the path stored from the analysis step
            presentation = powerpoint.Presentations.Open(os.path.abspath(self.current_presentation))

            # Check if presentation was opened successfully
            if presentation is None:
                messagebox.showerror("Lỗi", "Không thể mở file PowerPoint.")
                # Clean up powerpoint object if it was created
                if powerpoint is not None:
                    try:
                        powerpoint.Quit()
                    except:
                        pass
                return # Exit the function if opening failed

            # Start slideshow
            slideshow = presentation.SlideShowSettings.Run()
            slideshow_window = slideshow.View

            # Split the script by slide markers
            # Assuming script format is: KỊCH BẢN THUYẾT TRÌNH
            # --- Slide 1 ---
            # Content for slide 1
            # --- Slide 2 ---
            # Content for slide 2
            # ...
            script_lines = script.split('\n')
            slide_scripts = {}
            current_slide_num = None
            current_script_lines = []

            for line in script_lines:
                slide_match = re.match(r'^---\s*Slide\s*(\d+)\s*---', line)
                if slide_match:
                    # If we were processing a slide, store its script before starting the next
                    if current_slide_num is not None:
                        slide_scripts[current_slide_num] = "\n".join(current_script_lines).strip()
                    # Start processing a new slide
                    current_slide_num = int(slide_match.group(1))
                    current_script_lines = []
                elif current_slide_num is not None:
                    # Add line to the current slide's script
                    current_script_lines.append(line)

            # Store the last slide's script
            if current_slide_num is not None:
                slide_scripts[current_slide_num] = "\n".join(current_script_lines).strip()

            # Play introduction
            intro_text = "Xin chào, tôi là AI của công ty RATRACO. Sau đây tôi xin thuyết trình về nội dung của bài thuyết trình."
            if hasattr(self, 'user_voice_path') and self.user_voice_path and os.path.exists(self.user_voice_path):
                 pygame.mixer.music.load(self.user_voice_path)
                 pygame.mixer.music.play()
            else:
                # Use gTTS for intro if no user voice is recorded or file not found
                try:
                    tts = gTTS(text=intro_text, lang='vi', slow=False)
                    fp = BytesIO()
                    tts.write_to_fp(fp)
                    fp.seek(0)
                    pygame.mixer.music.load(fp)
                    pygame.mixer.music.play()
                except Exception as tts_intro_e:
                    print(f"Error generating intro speech: {tts_intro_e}")
                    # Continue without intro speech if TTS fails

            # Wait for intro to finish
            while pygame.mixer.music.get_busy():
                if is_presentation_paused:
                    pygame.mixer.music.pause()
                    while is_presentation_paused:
                        time.sleep(0.1)
                        self.root.update()
                    pygame.mixer.music.unpause() # This might not unpause if the mixer was manually paused
                time.sleep(0.1)
                self.root.update()

            # Ensure presentation starts from the first slide
            try:
                slideshow_window.GotoSlide(1)
            except Exception as goto_e:
                 print(f"Could not go to first slide: {goto_e}")
                 # If GoToSlide fails, assume it's already on the first slide or proceed

            # Process each slide using the script
            # Iterate through the actual slides in the slideshow window
            # Need to get the slide count from the running presentation
            total_slides_running = presentation.Slides.Count

            for i in range(1, total_slides_running + 1): # Slides are 1-indexed in win32com
                if is_presentation_paused:
                    while is_presentation_paused:
                        time.sleep(0.1)
                        self.root.update()

                # Update progress
                self.progress["value"] = int((i / total_slides_running) * 100)
                self.status_label.config(text=f"🗣️ Đang thuyết trình slide {i}/{total_slides_running}...")
                self.root.update_idletasks()

                # Get the script for the current slide number
                # Use .get() with a default empty string to avoid KeyError if a slide number is missing in the script
                slide_script_to_speak = slide_scripts.get(i, "(Không có nội dung kịch bản cho slide này)")

                # Speak the script for the current slide
                if slide_script_to_speak.strip():
                    try:
                        # Use gTTS for speaking the script
                        tts = gTTS(text=slide_script_to_speak, lang='vi', slow=False)
                        fp = BytesIO()
                        tts.write_to_fp(fp)
                        fp.seek(0)
                        # Stop any currently playing audio before loading the new one
                        pygame.mixer.music.stop()
                        pygame.mixer.music.load(fp)
                        pygame.mixer.music.play()

                        # Wait for speech to finish
                        while pygame.mixer.music.get_busy():
                            if is_presentation_paused:
                                pygame.mixer.music.pause()
                                while is_presentation_paused:
                                    time.sleep(0.1)
                                    self.root.update()
                                pygame.mixer.music.unpause()
                            time.sleep(0.1)
                            self.root.update()
                    except Exception as tts_slide_e:
                        print(f"Error generating speech for slide {i}: {tts_slide_e}")
                        # Continue to the next slide even if speech fails for this one

                # Wait a moment before advancing (optional, provides a brief pause)
                time.sleep(0.5)

                # Go to next slide if not the last one
                if i < total_slides_running:
                     # slideshow_window.Next() advances to the next slide
                    try:
                        slideshow_window.Next()
                    except Exception as next_slide_e:
                        print(f"Could not advance to next slide after slide {i}: {next_slide_e}")
                        # If advancing fails, might be the end or another issue, break the loop
                        break
                else:
                    # If it's the last slide, wait a bit longer before conclusion or ending
                    time.sleep(1)

            # Play conclusion
            conclusion_text = "Cảm ơn các bạn đã lắng nghe AI RATRACO."
            # Use user voice for conclusion if available, otherwise use gTTS
            if hasattr(self, 'user_voice_path') and self.user_voice_path and os.path.exists(self.user_voice_path):
                # Stop any remaining slide audio before playing conclusion
                pygame.mixer.music.stop()
                pygame.mixer.music.load(self.user_voice_path)
                pygame.mixer.music.play()
            else:
                try:
                    # Stop any remaining slide audio before playing conclusion
                    pygame.mixer.music.stop()
                    tts = gTTS(text=conclusion_text, lang='vi', slow=False)
                    fp = BytesIO()
                    tts.write_to_fp(fp)
                    fp.seek(0)
                    pygame.mixer.music.load(fp)
                    pygame.mixer.music.play()
                except Exception as tts_conc_e:
                     print(f"Error generating conclusion speech: {tts_conc_e}")
                     # Continue to cleanup if conclusion speech fails

            # Wait for conclusion to finish
            while pygame.mixer.music.get_busy():
                if is_presentation_paused:
                    pygame.mixer.music.pause()
                    while is_presentation_paused:
                        time.sleep(0.1)
                        self.root.update()
                    # No need to unpause the mixer here, the main loop handles it
                time.sleep(0.1)
                self.root.update()

            # Presentation finished, try to exit the slideshow
            try:
                 if slideshow_window is not None:
                     slideshow_window.Exit()
            except Exception as exit_e:
                 print(f"Error exiting slideshow: {exit_e}")

            # Clean up PowerPoint objects
            if presentation is not None:
                try:
                    presentation.Close()
                except:
                    pass
            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except:
                    pass

            # Disable pause button
            self.pause_button.config(state='disabled')

            self.status_label.config(text="✅ Đã hoàn thành thuyết trình!")
            self.progress["value"] = 100

        except Exception as e:
            # Catch any unexpected errors during the presentation process
            messagebox.showerror("Lỗi", f"Lỗi trong quá trình thuyết trình: {str(e)}")
            self.status_label.config(text="❌ Lỗi thuyết trình.")
            self.pause_button.config(state='disabled')

            # Clean up - ensure win32com objects are closed if they were somehow created here
            try:
                # Check if powerpoint_app was created and is still running
                if 'presentation' in locals() and presentation is not None:
                    # presentation here would be a win32com object if created
                    presentation.Close()
                if 'powerpoint' in locals() and powerpoint is not None:
                    # powerpoint here would be a win32com object if created
                    powerpoint.Quit()
            except Exception as cleanup_e:
                print(f"Error during error cleanup: {cleanup_e}")

    def toggle_presentation(self):
        """Toggle presentation pause/play state"""
        global is_presentation_paused
        is_presentation_paused = not is_presentation_paused
        if is_presentation_paused:
            self.pause_button.config(text="▶️ Tiếp tục")
            pygame.mixer.pause()
        else:
            self.pause_button.config(text="⏸️ Tạm dừng")
            pygame.mixer.unpause()

    def toggle_voice_recording(self):
        """Toggle voice recording for presentation"""
        global is_recording_voice, voice_recording_data, user_voice_path
        
        if not is_recording_voice:
            # Start recording
            if not self.mic_var.get():
                messagebox.showwarning("Chưa chọn micro", "Vui lòng chọn thiết bị micro trước khi ghi âm.")
                return
                
            try:
                device_index = int(self.mic_var.get().split(":")[0])
                is_recording_voice = True
                voice_recording_data = []
                
                self.stream = sd.InputStream(
                    device=device_index,
                    channels=1,
                    samplerate=SAMPLE_RATE,
                    callback=self.voice_callback
                )
                self.stream.start()
                
                self.voice_button.config(text="⏹️ Dừng ghi giọng")
                self.status_label.config(text="🎙️ Đang ghi giọng nói...")
                self.progress["value"] = 0
                
            except Exception as e:
                messagebox.showerror("Lỗi ghi âm", str(e))
        else:
            # Stop recording
            is_recording_voice = False
            try:
                if self.stream:
                    self.stream.stop()
                    self.stream.close()
            except:
                pass
                
            self.voice_button.config(text="🎙️ Ghi giọng nói")
            self.status_label.config(text="🛠️ Đang xử lý giọng nói...")
            
            # Process the recorded voice
            threading.Thread(target=self.process_voice_recording, daemon=True).start()

    def voice_callback(self, indata, frames, time, status):
        """Callback function for voice recording"""
        global voice_recording_data, is_recording_voice
        if is_recording_voice:
            volume = np.linalg.norm(indata) * 10
            volume = min(volume, 100)
            self.progress["value"] = volume
            voice_recording_data.append(indata.copy())

    def process_voice_recording(self):
        """Process the recorded voice data"""
        global voice_recording_data, user_voice_path, SAMPLE_RATE
        
        try:
            if not voice_recording_data:
                self.status_label.config(text="⚠️ Không có dữ liệu ghi âm")
                return
                
            audio_data = np.concatenate(voice_recording_data, axis=0)
            int16_audio = np.int16(audio_data * 32767)
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_file:
                write(temp_file.name, SAMPLE_RATE, int16_audio)
                user_voice_path = temp_file.name
            
            self.status_label.config(text="✅ Đã lưu giọng nói!")
            messagebox.showinfo("Thành công", "Đã lưu giọng nói thành công!")
            self.progress["value"] = 100
            
        except Exception as e:
            messagebox.showerror("Lỗi xử lý", str(e))
            self.status_label.config(text="❌ Lỗi xử lý giọng nói.")

    def analyze_text_emphasis(self, text):
        """Analyze text to determine emphasis points and speech patterns"""
        try:
            # Keywords that indicate importance
            emphasis_keywords = {
                'quan trọng': 1.5,
                'chính': 1.4,
                'đặc biệt': 1.4,
                'then chốt': 1.5,
                'mấu chốt': 1.5,
                'đáng chú ý': 1.3,
                'đáng kể': 1.3,
                'đáng quan tâm': 1.3,
                'tăng': 1.2,
                'giảm': 1.2,
                'cao nhất': 1.4,
                'thấp nhất': 1.4,
                'đầu tiên': 1.2,
                'cuối cùng': 1.2,
                'kết luận': 1.3,
                'tóm lại': 1.3
            }
            
            # Split text into sentences
            sentences = re.split(r'[.!?]+', text)
            sentences = [s.strip() for s in sentences if s.strip()]
            
            # Analyze each sentence
            analyzed_sentences = []
            for sentence in sentences:
                emphasis_level = 1.0
                emphasis_points = []
                
                # Check for emphasis keywords
                for keyword, level in emphasis_keywords.items():
                    if keyword in sentence.lower():
                        emphasis_level = max(emphasis_level, level)
                        emphasis_points.append(keyword)
                
                # Check for numbers and percentages
                numbers = re.findall(r'\d+(?:[.,]\d+)?%?', sentence)
                if numbers:
                    emphasis_level = max(emphasis_level, 1.3)
                    emphasis_points.extend(numbers)
                
                # Check for question marks
                if '?' in sentence:
                    emphasis_level = max(emphasis_level, 1.4)
                    emphasis_points.append('câu hỏi')
                
                # Check for exclamation marks
                if '!' in sentence:
                    emphasis_level = max(emphasis_level, 1.5)
                    emphasis_points.append('cảm thán')
                
                analyzed_sentences.append({
                    'text': sentence,
                    'emphasis_level': emphasis_level,
                    'emphasis_points': emphasis_points
                })
            
            return analyzed_sentences
            
        except Exception as e:
            print(f"Error analyzing text emphasis: {str(e)}")
            return [{'text': text, 'emphasis_level': 1.0, 'emphasis_points': []}]

    def process_slide_content(self, slide_text, slide_image):
        """Process slide content in detail"""
        try:
            # Initialize detailed analysis
            analysis = {
                'text_content': [],
                'charts': [],
                'images': [],
                'tables': [],
                'key_points': [],
                'emphasis_points': []
            }
            
            # Analyze text content
            text_analysis = self.analyze_text_emphasis(slide_text)
            analysis['text_content'] = text_analysis
            
            # Extract key points
            for sentence in text_analysis:
                if sentence['emphasis_level'] >= 1.3:
                    analysis['key_points'].append(sentence['text'])
                if sentence['emphasis_points']:
                    analysis['emphasis_points'].extend(sentence['emphasis_points'])
            
            # Analyze visual elements if slide image is provided
            if slide_image:
                visual_analysis = analyze_slide_content(slide_image)
                if visual_analysis:
                    analysis.update(visual_analysis)
            
            return analysis
            
        except Exception as e:
            print(f"Error processing slide content: {str(e)}")
            return None

    def generate_detailed_script(self, analysis):
        """Generate detailed presentation script from analysis"""
        try:
            script = []
            
            # Process text content
            for sentence in analysis['text_content']:
                text = sentence['text']
                script.append(text)
            
            # Add chart descriptions
            for chart in analysis['charts']:
                script.append(f"\nBiểu đồ {chart['type']}:")
                script.append(chart['description'])
                if chart['data']:
                    script.append(f"Dữ liệu: {', '.join(map(str, chart['data']))}")
            
            # Add image descriptions
            for image in analysis['images']:
                script.append(f"\nHình ảnh: {image['description']}")
            
            # Add table descriptions
            for table in analysis['tables']:
                script.append(f"\nBảng: {table['description']}")
                if table['data']:
                    script.append("Dữ liệu chi tiết:")
                    for row in table['data']:
                        script.append(" | ".join(row))
            
            return "\n".join(script)
            
        except Exception as e:
            print(f"Error generating detailed script: {str(e)}")
            return ""


if __name__ == "__main__":
    root = tk.Tk()
    app = RatracoAITool(root)
    root.mainloop()
























#Author: kiniemboquenjerry@gmail.com